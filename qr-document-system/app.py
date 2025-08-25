from flask import Flask, render_template, request, jsonify, send_file, url_for
import qrcode
import os
from werkzeug.utils import secure_filename
import pandas as pd
from docx import Document
from pptx import Presentation
import io
import base64
from PIL import Image
import uuid
import json
from datetime import datetime

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['IMAGES_FOLDER'] = 'static/images'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# 确保文件夹存在
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs('static', exist_ok=True)
os.makedirs(app.config['IMAGES_FOLDER'], exist_ok=True)

# 允许的文件扩展名
ALLOWED_EXTENSIONS = {'docx', 'pptx', 'xlsx', 'xls'}

# 存储已处理的文件内容
document_storage = {}

# 当前会话ID（用于生成唯一的内容展示页面）
current_session = {
    'session_id': None,
    'documents': {},
    'created_at': None
}


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def generate_qr_code(session_id=None):
    """生成二维码"""
    if session_id:
        website_url = f"http://192.168.1.110:8080/view/{session_id}"
        qr_filename = f'static/qrcode_{session_id}.png'
    else:
        website_url = "http://192.168.1.110:8080"
        qr_filename = 'static/qrcode.png'

    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=4,
    )
    qr.add_data(website_url)
    qr.make(fit=True)

    img = qr.make_image(fill_color="black", back_color="white")
    img.save(qr_filename)
    print(f"二维码已生成并保存到 {qr_filename}")
    print(f"二维码指向: {website_url}")

    return qr_filename, website_url


def read_word_file(filepath):
    """读取Word文件内容"""
    try:
        doc = Document(filepath)
        content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content.append(paragraph.text)
        return '\n'.join(content)
    except Exception as e:
        return f"读取Word文件出错: {str(e)}"


def read_excel_file(filepath):
    """读取Excel文件内容"""
    try:
        excel_file = pd.ExcelFile(filepath)
        sheets_data = {}

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(filepath, sheet_name=sheet_name)
            sheets_data[sheet_name] = df.to_html(classes='table table-striped', index=False)

        return sheets_data
    except Exception as e:
        return {"error": f"读取Excel文件出错: {str(e)}"}


def extract_images_from_ppt(presentation, session_id, slide_index):
    """从PPT幻灯片中提取图片"""
    images = []
    slide = presentation.slides[slide_index]

    try:
        for shape_index, shape in enumerate(slide.shapes):
            # 检查是否有图片
            if hasattr(shape, 'image'):
                try:
                    image = shape.image
                    image_bytes = image.blob

                    # 生成唯一文件名
                    image_filename = f"slide_{slide_index + 1}_img_{shape_index + 1}_{session_id}.png"
                    image_path = os.path.join(app.config['IMAGES_FOLDER'], image_filename)

                    # 保存图片
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)

                    images.append({
                        'filename': image_filename,
                        'path': f"images/{image_filename}",
                        'alt': f"幻灯片{slide_index + 1}图片{shape_index + 1}"
                    })

                except Exception as e:
                    print(f"提取图片出错: {e}")
                    continue
    except Exception as e:
        print(f"处理幻灯片图片时出错: {e}")

    return images


def read_ppt_file(filepath, session_id=None):
    """读取PPT文件内容"""
    try:
        presentation = Presentation(filepath)
        slides_content = []

        for i, slide in enumerate(presentation.slides):
            slide_content = {
                'slide_number': i + 1,
                'content': [],
                'images': []
            }

            # 如果有session_id，提取图片
            if session_id:
                images = extract_images_from_ppt(presentation, session_id, i)
                slide_content['images'] = images

            # 提取文字内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content['content'].append(shape.text.strip())
                elif shape.has_table:
                    table_text = []
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        table_text.append(" | ".join(row_text))
                    if table_text:
                        slide_content['content'].append("表格内容:\n" + "\n".join(table_text))

            slides_content.append(slide_content)

        return slides_content

    except Exception as e:
        return [{"error": f"读取PPT文件出错: {str(e)}"}]


@app.route('/')
def index():
    """主页"""
    return render_template('index.html')


@app.route('/upload', methods=['POST'])
def upload_file():
    """文件上传处理"""
    global current_session

    if 'file' not in request.files:
        return jsonify({'error': '没有文件被上传'})

    file = request.files['file']
    file_type = request.form.get('file_type')

    if file.filename == '':
        return jsonify({'error': '没有选择文件'})

    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)

        # 如果是新的会话，创建会话ID
        if not current_session['session_id']:
            current_session['session_id'] = str(uuid.uuid4())[:8]
            current_session['created_at'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            current_session['documents'] = {}

        # 根据文件类型处理文件
        content = None
        content_type = None

        if file_type == 'word':
            content = read_word_file(filepath)
            content_type = 'text'
        elif file_type == 'excel':
            content = read_excel_file(filepath)
            content_type = 'table'
        elif file_type == 'ppt':
            content = read_ppt_file(filepath, current_session['session_id'])
            content_type = 'slides'
        else:
            return jsonify({'error': '不支持的文件类型'})

        # 保存到当前会话
        current_session['documents'][file_type] = {
            'filename': filename,
            'content': content,
            'type': content_type,
            'uploaded_at': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }

        # 将会话保存到全局存储
        document_storage[current_session['session_id']] = current_session.copy()

        # 生成新的二维码
        qr_filename, qr_url = generate_qr_code(current_session['session_id'])

        return jsonify({
            'success': True,
            'content': content,
            'type': content_type,
            'session_id': current_session['session_id'],
            'qr_url': qr_url,
            'qr_filename': qr_filename.replace('static/', '')
        })

    return jsonify({'error': '文件格式不支持'})


@app.route('/qr')
def show_qr():
    """显示二维码页面"""
    return render_template('qr.html')


@app.route('/view/<session_id>')
def view_documents(session_id):
    """查看已上传的文档内容"""
    if session_id not in document_storage:
        return render_template('error.html', error="会话不存在或已过期")

    session_data = document_storage[session_id]
    return render_template('view.html', session_data=session_data)


@app.route('/view/<session_id>/<doc_type>')
def view_document_detail(session_id, doc_type):
    """查看单个文档的详细内容"""
    if session_id not in document_storage:
        return render_template('error.html', error="会话不存在或已过期")

    session_data = document_storage[session_id]

    if doc_type not in session_data['documents']:
        return render_template('error.html', error=f"找不到{doc_type}类型的文档")

    doc_data = session_data['documents'][doc_type]

    # 如果是PPT，跳转到增强展示页面
    if doc_type == 'ppt':
        return render_template('ppt_enhanced.html',
                               session_data=session_data,
                               ppt_data=doc_data)

    return render_template('document_detail.html',
                           session_data=session_data,
                           doc_type=doc_type,
                           doc_data=doc_data)


@app.route('/new_session')
def new_session():
    """创建新会话"""
    global current_session
    current_session = {
        'session_id': None,
        'documents': {},
        'created_at': None
    }
    return jsonify({'success': True, 'message': '新会话已创建'})


@app.route('/current_qr')
def get_current_qr():
    """获取当前会话的二维码信息"""
    if current_session['session_id']:
        qr_filename = f"qrcode_{current_session['session_id']}.png"
        qr_url = f"http://192.168.1.110:8080/view/{current_session['session_id']}"
        return jsonify({
            'has_qr': True,
            'qr_filename': qr_filename,
            'qr_url': qr_url,
            'session_id': current_session['session_id'],
            'document_count': len(current_session['documents'])
        })
    else:
        return jsonify({'has_qr': False})


if __name__ == '__main__':
    # 生成默认二维码
    generate_qr_code()

    print("启动Flask应用...")
    print("访问 http://192.168.1.110:8080 查看网站")
    print("访问 http://192.168.1.110:8080/qr 查看二维码")

    # 启动Flask应用
    app.run(debug=True, host='0.0.0.0', port=8080)